"""Genera la presentacion profesional en PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colores ---
BG_DARK = RGBColor(0x0d, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1b, 0x22)
BG_CARD2 = RGBColor(0x1e, 0x2a, 0x3a)
BLUE = RGBColor(0x34, 0x98, 0xdb)
GREEN = RGBColor(0x27, 0xae, 0x60)
RED = RGBColor(0xe7, 0x4c, 0x3c)
ORANGE = RGBColor(0xf3, 0x9c, 0x12)
PURPLE = RGBColor(0x9b, 0x59, 0xb6)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xe0, 0xe0, 0xe0)
GREY = RGBColor(0x6c, 0x80, 0x91)
DARK_GREY = RGBColor(0x55, 0x66, 0x77)
BORDER = RGBColor(0x1a, 0x3a, 0x5c)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
FONT = "Calibri"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = alignment
    return tb


def add_bullet_box(slide, left, top, width, bullets, size=14, color=LIGHT):
    tb = slide.shapes.add_textbox(left, top, width, Inches(len(bullets) * 0.38))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_before = Pt(6)
        p.space_after = Pt(2)


def add_kpi_card(slide, left, top, width, height, value, label, accent_color=BLUE):
    add_shape(slide, left, top, width, height, BG_CARD2)
    add_rect(slide, left, top, width, Pt(3), accent_color)
    add_text(slide, left, top + Inches(0.2), width, Inches(0.5),
             value, size=26, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.7), width, Inches(0.35),
             label, size=11, color=GREY, alignment=PP_ALIGN.CENTER)


def add_transition(slide, text, y=6.95):
    add_text(slide, Inches(0.8), Inches(y), Inches(11.7), Inches(0.3),
             text, size=13, color=DARK_GREY, alignment=PP_ALIGN.CENTER)


def add_flow_step(slide, x, y, w, h, text, color, size=11):
    add_shape(slide, x, y, w, h, color)
    add_text(slide, x, y + Pt(2), w, h,
             text, size=size, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def add_flow_arrow(slide, x, y):
    add_text(slide, x, y, Inches(0.4), Inches(0.4),
             "→", size=18, color=GREY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1 — PORTADA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_rect(slide, 0, 0, prs.slide_width, Pt(4), BLUE)

add_text(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.5),
         "Prediccion de Fuga de Clientes\nen el Sector de Automocion",
         size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.6),
         "Machine Learning + Estrategia Comercial Rentable",
         size=22, color=BLUE, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.5), Inches(4.6), Inches(4.3), Pt(1), DARK_GREY)

add_text(slide, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.5),
         "Caso Practico IA  ·  Universidad Alfonso X el Sabio  ·  2026",
         size=14, color=GREY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.5),
         "Python  ·  scikit-learn  ·  SHAP  ·  Plotly  ·  Streamlit",
         size=12, color=DARK_GREY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — PROBLEMA + DATOS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.3),
         "01  EL PROBLEMA", size=11, color=BLUE, bold=True)

add_text(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.5),
         "Perder un cliente cuesta 5-7x mas que retenerlo",
         size=28, color=WHITE, bold=True)

# Left card
add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(3.5), BG_CARD2, BORDER)
add_text(slide, Inches(1.1), Inches(1.5), Inches(5), Inches(0.35),
         "El Problema", size=15, color=BLUE, bold=True)
add_bullet_box(slide, Inches(1.1), Inches(1.9), Inches(5), [
    "Un cliente en fuga lleva >400 dias sin revision en taller",
    "Sin intervencion, se pierde el ingreso de mantenimiento",
    "Cada cliente genera 2.038 EUR de margen en 5 revisiones",
    "Captar un cliente nuevo cuesta 5-7x mas que retenerlo",
], size=13)

# Right card
add_shape(slide, Inches(6.9), Inches(1.4), Inches(5.6), Inches(3.5), BG_CARD2, BORDER)
add_text(slide, Inches(7.2), Inches(1.5), Inches(5), Inches(0.35),
         "Los Datos", size=15, color=BLUE, bold=True)
add_bullet_box(slide, Inches(7.2), Inches(1.9), Inches(5), [
    "58.049 registros  /  44.053 clientes unicos",
    "8,8% tasa de churn (ratio 10:1, desbalanceado)",
    "27.070 con 0 revisiones excluidos del entrenamiento",
    "10.000 clientes nuevos en produccion (sin etiqueta)",
], size=13)

# Revenue at Risk banner
add_shape(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(0.95), BG_CARD, RED)
add_text(slide, Inches(2.5), Inches(5.3), Inches(8.3), Inches(0.45),
         "4,9M EUR en Revenue at Risk", size=26, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2.5), Inches(5.75), Inches(8.3), Inches(0.3),
         "en la cartera de 10.000 clientes de produccion",
         size=12, color=GREY, alignment=PP_ALIGN.CENTER)

add_transition(slide, "¿Podemos predecir quienes se van a ir antes de que ocurra?", y=6.5)


# ============================================================
# SLIDE 3 — EL MODELO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.3),
         "02  EL MODELO", size=11, color=BLUE, bold=True)

add_text(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.5),
         "Random Forest: detectamos al 71% de los clientes en riesgo",
         size=28, color=WHITE, bold=True)

# Pipeline
add_shape(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.85), BG_CARD2)
steps = ["Limpieza", "EDA", "FE (64 vars)", "4 Algoritmos", "Random Forest"]
for i, step in enumerate(steps):
    x = Inches(1.1) + i * Inches(2.3)
    bg = BLUE if i == 4 else BORDER
    add_shape(slide, x, Inches(1.55), Inches(1.7), Inches(0.5), bg)
    add_text(slide, x, Inches(1.6), Inches(1.7), Inches(0.4),
             step, size=12, color=WHITE, bold=(i == 4), alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, x + Inches(1.7), Inches(1.55), Inches(0.6), Inches(0.45),
                 "→", size=18, color=GREY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(2.35), Inches(11.7), Inches(0.25),
         "Variables de revisiones y temporales excluidas para garantizar generalizacion a clientes nuevos",
         size=11, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

# KPIs
kpis = [
    ("0,789", "AUC-ROC", BLUE),
    ("71%", "Recall", GREEN),
    ("34,2%", "Precision", ORANGE),
    ("0,46", "F1-Score", PURPLE),
    ("64", "Features", GREY),
]
for i, (val, label, col) in enumerate(kpis):
    x = Inches(0.8) + i * Inches(2.45)
    add_kpi_card(slide, x, Inches(2.85), Inches(2.1), Inches(1.05), val, label, col)

# Left: Top Features
add_text(slide, Inches(0.8), Inches(4.15), Inches(4), Inches(0.35),
         "Top Features", size=15, color=BLUE, bold=True)

features = [
    ("EN_GARANTIA", "fuera de garantia = mayor riesgo", 0.95),
    ("Modelo vehiculo", "H=0% churn, C/A >13%", 0.72),
    ("QUEJA", "queja incrementa la probabilidad", 0.55),
    ("Seguro bateria", "sin seguro = mayor riesgo", 0.42),
    ("Forma de pago", "contado > financiacion", 0.35),
]
for i, (name, desc, imp) in enumerate(features):
    y = Inches(4.6) + i * Inches(0.42)
    add_text(slide, Inches(0.8), y, Inches(1.8), Inches(0.35),
             name, size=11, color=LIGHT, bold=True)
    bw = int(Inches(3.0) * imp)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), y + Pt(4), bw, Pt(14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    bar.adjustments[0] = 0.3
    add_text(slide, Inches(6.0), y, Inches(3), Inches(0.35),
             desc, size=11, color=GREY)

# Right: Precision note
add_shape(slide, Inches(9.3), Inches(4.15), Inches(3.2), Inches(1.4), BG_CARD2, BORDER)
add_text(slide, Inches(9.5), Inches(4.25), Inches(2.8), Inches(0.3),
         "Trade-off Precision/Recall", size=12, color=BLUE, bold=True)
add_text(slide, Inches(9.5), Inches(4.6), Inches(2.8), Inches(0.8),
         "Precision 34% = por cada 3 alertas,\n1 es churn real.\n\nPreferimos contactar de mas\na perder un cliente.",
         size=11, color=LIGHT)

add_transition(slide, "Sabemos quien se va. Pero ¿que hacemos con eso?")


# ============================================================
# SLIDE 4 — ESTRATEGIA COMERCIAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.3),
         "03  ESTRATEGIA COMERCIAL", size=11, color=BLUE, bold=True)

add_text(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.5),
         "4 acciones automaticas. ROI de 2,2x. Margen minimo garantizado.",
         size=28, color=WHITE, bold=True)

# --- Decision flow diagram ---
add_shape(slide, Inches(0.8), Inches(1.35), Inches(11.7), Inches(1.1), BG_CARD2)
add_text(slide, Inches(1.0), Inches(1.38), Inches(3), Inches(0.25),
         "Motor de decision", size=10, color=GREY, bold=True)

# Flow: Cliente → Modelo RF → Prob Churn → Riesgo + CLTV → Accion → ROI
flow = [
    ("Cliente", BORDER, 1.0),
    ("Modelo RF", BORDER, 3.0),
    ("Prob. Churn", BORDER, 5.0),
    ("Riesgo + CLTV", BLUE, 7.0),
    ("Accion", GREEN, 9.2),
    ("ROI", ORANGE, 11.1),
]
for label, col, x in flow:
    add_flow_step(slide, Inches(x), Inches(1.7), Inches(1.5), Inches(0.45), label, col, size=11)

for _, _, x in flow[:-1]:
    add_flow_arrow(slide, Inches(x) + Inches(1.5), Inches(1.7))

# Table card
add_shape(slide, Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.85), BG_CARD2, BORDER)

# Headers
hdrs = ["Accion", "Regla", "Descuento", "Contacto", "Conv.", "Clientes", "ROI"]
hx = [1.1, 3.2, 5.8, 7.2, 8.5, 9.6, 10.8]
hw = [2.0, 2.5, 1.3, 1.2, 1.0, 1.1, 1.0]
for h, x, w in zip(hdrs, hx, hw):
    add_text(slide, Inches(x), Inches(2.8), Inches(w), Inches(0.35),
             h, size=11, color=GREY, bold=True)

add_rect(slide, Inches(1.1), Inches(3.15), Inches(11.1), Pt(1), BORDER)

# Rows
rows = [
    (RED, "Ret. Premium", "Alto + CLTV alto", "30% C(n)", "50 EUR", "42%", "204", "0,4x"),
    (ORANGE, "Ret. Estandar", "Alto + CLTV bajo", "22% C(n)", "35 EUR", "40%", "1.495", "0,7x"),
    (BLUE, "Fidelizacion", "Riesgo medio", "15% C(n)", "20 EUR", "38%", "2.060", "1,3x"),
    (GREEN, "Mant. Preventivo", "Riesgo bajo", "8% C(n)", "10 EUR", "50%", "6.241", "3,8x"),
]
for i, (color, *vals) in enumerate(rows):
    y = Inches(3.25) + i * Inches(0.55)
    add_rect(slide, Inches(0.85), y + Pt(4), Pt(4), Inches(0.35), color)
    for val, x, w in zip(vals, hx, hw):
        add_text(slide, Inches(x), y, Inches(w), Inches(0.45),
                 val, size=12, color=LIGHT)

# Formula card (bottom left)
add_shape(slide, Inches(0.8), Inches(5.8), Inches(6.0), Inches(1.4), BG_CARD2)
add_text(slide, Inches(1.1), Inches(5.9), Inches(5.5), Inches(0.3),
         "Formula ROI ajustado", size=13, color=BLUE, bold=True)
add_text(slide, Inches(1.1), Inches(6.2), Inches(5.5), Inches(0.35),
         "ROI = (Margen neto x Tasa conversion) / Inversion",
         size=15, color=WHITE, bold=True)
add_text(slide, Inches(1.1), Inches(6.6), Inches(5.5), Inches(0.45),
         "No todos los clientes contactados vuelven. Los de menor riesgo\nresponden mejor (50%) que los de alto riesgo (42%).",
         size=11, color=GREY)

# 3 KPIs (bottom right)
kpi_data = [
    ("547K EUR", "Inversion total", ORANGE),
    ("2,2x", "ROI global", GREEN),
    ("37%", "Margen minimo", BLUE),
]
for i, (val, label, col) in enumerate(kpi_data):
    x = Inches(7.2) + i * Inches(1.85)
    add_kpi_card(slide, x, Inches(5.8), Inches(1.6), Inches(1.05), val, label, col)


# ============================================================
# SLIDE 5 — DASHBOARD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.3),
         "04  PRODUCTO FINAL", size=11, color=BLUE, bold=True)

add_text(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.5),
         "Dashboard interactivo desplegado en produccion",
         size=28, color=WHITE, bold=True)

# 6 Tab cards (2 rows x 3 cols)
tabs = [
    ("1", "Resumen", "KPIs globales, riesgo, churn por modelo", BLUE),
    ("2", "Explorador", "Ficha cliente + SHAP + What-If (11 vars)", GREEN),
    ("3", "Estrategia", "ROI por accion, CLTV, comparador A/B", ORANGE),
    ("4", "Sensibilidad", "Threshold, ajuste descuentos en vivo", PURPLE),
    ("5", "Piloto Auto", "Sube CSV → predice → descarga Excel", RED),
    ("6", "Modelo", "ROC, PR, confusion matrix, features", GREY),
]
for i, (num, name, desc, col) in enumerate(tabs):
    row = i // 3
    col_idx = i % 3
    x = Inches(0.8) + col_idx * Inches(4.15)
    y = Inches(1.4) + row * Inches(1.55)

    add_shape(slide, x, y, Inches(3.85), Inches(1.25), BG_CARD2)
    add_rect(slide, x, y, Inches(3.85), Pt(3), col)
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.35),
             num, size=20, color=col, bold=True)
    add_text(slide, x + Inches(0.5), y + Inches(0.15), Inches(3.1), Inches(0.35),
             name, size=15, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.5), y + Inches(0.55), Inches(3.1), Inches(0.55),
             desc, size=12, color=GREY)

# URL banner
add_shape(slide, Inches(2.0), Inches(4.7), Inches(9.3), Inches(0.85), BG_CARD, BLUE)
add_text(slide, Inches(2.0), Inches(4.8), Inches(9.3), Inches(0.35),
         "churn-prediction-automotive-jhxywrrmnkg5nfba3vxsxk.streamlit.app",
         size=15, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2.0), Inches(5.15), Inches(9.3), Inches(0.3),
         "Demo en vivo — accesible desde cualquier navegador",
         size=11, color=GREY, alignment=PP_ALIGN.CENTER)

# Stack
add_text(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.35),
         "Python 3.11  ·  scikit-learn  ·  SHAP  ·  Plotly  ·  Streamlit  ·  GitHub",
         size=13, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

# Note
add_shape(slide, Inches(1.5), Inches(6.35), Inches(10.3), Inches(0.55), BG_CARD2)
add_text(slide, Inches(1.7), Inches(6.4), Inches(9.9), Inches(0.4),
         "El dashboard convierte un modelo academico en un producto real que un equipo comercial podria usar manana.",
         size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_transition(slide, "Todo esto se traduce en resultados medibles.")


# ============================================================
# SLIDE 6 — CONCLUSIONES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.3),
         "05  CONCLUSIONES", size=11, color=BLUE, bold=True)

add_text(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.5),
         "Resultados e impacto", size=28, color=WHITE, bold=True)

# 4 Result cards (2x2)
results = [
    ("71% de los clientes en riesgo detectados",
     "Rendimiento solido para un dataset con ratio 10:1", GREEN),
    ("Variables predictivas accionables",
     "Garantia, modelo y quejas — factores que el concesionario controla", BLUE),
    ("ROI de 2,2x sobre 10.000 clientes",
     "547K EUR protegiendo 4,9M EUR en Revenue at Risk", ORANGE),
    ("Dashboard en produccion",
     "6 tabs con SHAP, What-If, sensibilidad y pipeline automatico", PURPLE),
]
for i, (title, desc, col) in enumerate(results):
    row = i // 2
    col_idx = i % 2
    x = Inches(0.8) + col_idx * Inches(6.15)
    y = Inches(1.4) + row * Inches(1.5)

    add_shape(slide, x, y, Inches(5.85), Inches(1.2), BG_CARD2)
    add_rect(slide, x, y, Pt(4), Inches(1.2), col)
    add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(5.3), Inches(0.35),
             title, size=15, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.55), Inches(5.3), Inches(0.5),
             desc, size=12, color=GREY)

# Limitations
add_text(slide, Inches(0.8), Inches(4.6), Inches(4), Inches(0.35),
         "Limitaciones y trabajo futuro", size=13, color=GREY, bold=True)

add_bullet_box(slide, Inches(0.8), Inches(4.95), Inches(11.7), [
    "Precision del 34% — mejorable con datos de CRM y comportamiento web",
    "Validacion A/B en produccion pendiente para medir impacto real",
], size=12, color=DARK_GREY)

# Hero statement
add_shape(slide, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.85), BG_CARD, GREEN)
add_text(slide, Inches(1.5), Inches(5.95), Inches(10.3), Inches(0.55),
         "Por cada euro invertido, recuperamos 2,2 EUR en margen de mantenimiento",
         size=22, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — GRACIAS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_rect(slide, 0, 0, prs.slide_width, Pt(4), BLUE)

add_text(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(0.8),
         "Gracias", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.5),
         "¿Preguntas?", size=22, color=GREY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.0), Inches(4.2), Inches(3.3), Pt(1), DARK_GREY)

add_text(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.4),
         "Dashboard: churn-prediction-automotive-jhxywrrmnkg5nfba3vxsxk.streamlit.app",
         size=14, color=BLUE, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.4),
         "GitHub: github.com/PabloJimenezGirardeau/churn-prediction-automotive",
         size=14, color=DARK_GREY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
out = os.path.join(os.path.dirname(__file__), "Presentacion_Churn_Prediction.pptx")
prs.save(out)
print(f"OK: {out}")
